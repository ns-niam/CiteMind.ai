"""
CiteMind AI - Presentation
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# Color palette
PURPLE = RGBColor(0x66, 0x7E, 0xEA)
DARK_PURPLE = RGBColor(0x76, 0x4B, 0xA2)
DARK_BG = RGBColor(0x0F, 0x0C, 0x29)
LIGHT_TEXT = RGBColor(0xF5, 0xF7, 0xFA)
ACCENT = RGBColor(0xA8, 0xED, 0xEA)
GRAY = RGBColor(0x94, 0xA3, 0xB8)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, DARK_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "CiteMind AI"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(72)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = ACCENT

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.3), Inches(1))
    sf = sub_box.text_frame
    sf.text = "Intelligent Research Assistant with Verifiable Citations"
    p = sf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(28)
    p.runs[0].font.color.rgb = LIGHT_TEXT
    p.runs[0].font.italic = True

    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.6))
    tg = tag_box.text_frame
    tg.text = "Powered by RAG · Groq · Gemini · ChromaDB"
    p = tg.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(18)
    p.runs[0].font.color.rgb = PURPLE

    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6))
    af = author_box.text_frame
    af.text = "Niam  ·  Machine Learning System Design  ·  May 2026"
    p = af.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(14)
    p.runs[0].font.color.rgb = GRAY


def add_content_slide(prs, title, bullets, footer_num=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)

    # Title bar
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.9))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.runs[0].font.size = Pt(36)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = ACCENT

    # Accent line
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.35), Inches(2.0), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = PURPLE
    line.line.fill.background()

    # Bullets
    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(12.0), Inches(5.5))
    bf = body.text_frame
    bf.word_wrap = True
    for i, item in enumerate(bullets):
        if i == 0:
            p = bf.paragraphs[0]
        else:
            p = bf.add_paragraph()
        p.text = "•  " + item
        p.runs[0].font.size = Pt(20)
        p.runs[0].font.color.rgb = LIGHT_TEXT
        p.space_after = Pt(12)

    # Footer
    if footer_num:
        ft = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.7), Inches(0.4))
        f = ft.text_frame
        f.text = f"{footer_num}"
        p = f.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.runs[0].font.size = Pt(12)
        p.runs[0].font.color.rgb = GRAY

    by = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(6.0), Inches(0.4))
    by_tf = by.text_frame
    by_tf.text = "CiteMind AI · by Niam"
    p = by_tf.paragraphs[0]
    p.runs[0].font.size = Pt(11)
    p.runs[0].font.color.rgb = GRAY


def add_image_slide(prs, title, image_path, caption=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.9))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.runs[0].font.size = Pt(32)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = ACCENT

    if Path(image_path).exists():
        slide.shapes.add_picture(image_path, Inches(2.0), Inches(1.6),
                                 width=Inches(9.0))

    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.6), Inches(12.0), Inches(0.5))
        cf = cap_box.text_frame
        cf.text = caption
        p = cf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(14)
        p.runs[0].font.italic = True
        p.runs[0].font.color.rgb = GRAY

    by = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(6.0), Inches(0.4))
    by_tf = by.text_frame
    by_tf.text = "CiteMind AI · by Niam"
    p = by_tf.paragraphs[0]
    p.runs[0].font.size = Pt(11)
    p.runs[0].font.color.rgb = GRAY


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_thank_you(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(2))
    tf = title_box.text_frame
    tf.text = "Thank You"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(80)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = ACCENT

    sub = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(12.3), Inches(1))
    sf = sub.text_frame
    sf.text = "Questions?"
    p = sf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(36)
    p.runs[0].font.color.rgb = LIGHT_TEXT

    repo = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.6))
    rf = repo.text_frame
    rf.text = "github.com/ns-niam/CiteMind.ai"
    p = rf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(20)
    p.runs[0].font.color.rgb = PURPLE


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title
    add_title_slide(prs)

    # 2. Problem Statement
    add_content_slide(prs, "1. Problem Statement", [
        "3M+ research papers published annually — researchers drown in information",
        "Keyword search returns documents, not synthesized answers",
        "LLMs hallucinate facts and fabricate citations — unsafe for academic work",
        "No verifiable sources → AI-generated summaries cannot be trusted",
        "PhD students spend 23+ hrs/week on literature review (Nature 2023)",
    ], footer_num="2/12")

    # 3. Actuality & Relevance
    add_content_slide(prs, "2. Actuality & Relevance", [
        "Educational: democratizes literature review for under-resourced students",
        "Medical: evidence-based decision support with verifiable citations",
        "Business / R&D: 10x faster technical due diligence",
        "Scientific: reduces time-to-insight when entering new fields",
        "Free LLMs (Groq, Gemini) + local vector DBs make this achievable now",
    ], footer_num="3/12")

    # 4. Novelty & Originality
    add_content_slide(prs, "3. Novelty & Originality", [
        "Dual-LLM Architecture: Groq (fast) + Gemini (deep) with comparison",
        "Confidence-Aware Refusal: 3-tier (High/Medium/Low) prevents hallucination",
        "Per-Chunk Citation Tracking: source + page + relevance score",
        "100% open-source, self-hostable, zero-cost (free LLM tiers)",
        "RAGAS-style evaluation built-in — quantitatively measurable quality",
    ], footer_num="4/12")

    # 5. Related Work
    add_content_slide(prs, "4. Related Work", [
        "Lewis et al. (2020) — RAG paradigm (NeurIPS 2020)",
        "Karpukhin et al. (2020) — Dense Passage Retrieval (EMNLP)",
        "Reimers & Gurevych (2019) — Sentence-BERT embeddings",
        "Es et al. (2023) — RAGAS evaluation framework (EACL 2024)",
        "Gao et al. (2023) — RAG survey identifies our gap: open-source + cited",
        "Compared to: ChatGPT (no cite), Perplexity (closed), NotebookLM (closed)",
    ], footer_num="5/12")

    # 6. Methodology - Architecture
    add_image_slide(
        prs, "5. Methodology — System Architecture",
        "assets/charts/architecture_diagram.png",
        caption="End-to-end RAG pipeline: Documents → Embed → Index → Retrieve → Generate → Cite"
    )

    # 7. Methodology details
    add_content_slide(prs, "5. Methodology — Components", [
        "Embedding: sentence-transformers/all-MiniLM-L6-v2 (384-dim, 80MB, CPU-fast)",
        "Vector DB: ChromaDB with HNSW index, cosine similarity",
        "Chunking: Recursive (1000 chars, 200 overlap) — respects sentence boundaries",
        "Retrieval: Top-K=5 with MMR re-ranking (λ=0.5) for diversity",
        "Generation: Groq Llama 3.3 70B + Gemini 2.5 Flash with citation prompts",
        "Confidence: top-1 score < 0.30 → refuse; 0.30–0.50 → caveat; > 0.50 → answer",
    ], footer_num="7/12")

    # 8. Results
    add_image_slide(
        prs, "6. Results — Metrics Comparison",
        "assets/charts/01_metrics_comparison.png",
        caption="Both LLMs exceed target on Faithfulness — RAG anti-hallucination works"
    )

    # 9. Results - radar
    add_image_slide(
        prs, "6. Results — Quality Profile",
        "assets/charts/05_radar_chart.png",
        caption="Multi-metric profile: Groq vs Gemini head-to-head"
    )

    # 10. Visualizations
    add_image_slide(
        prs, "7. Visualizations — Per-Query Faithfulness",
        "assets/charts/03_per_query_faithfulness.png",
        caption="10 hand-crafted eval queries, both LLMs above target line (0.85)"
    )

    # 11. Conclusions
    add_content_slide(prs, "8. Conclusions", [
        "Built a production-grade citation-aware research assistant — open-source",
        "All success criteria met: <5s latency, faithfulness > target, dual LLM",
        "Confidence layer turns 'trust me' into 'verify me' — defining contribution",
        "Reproducible: Jupyter notebook runs end-to-end, full eval suite included",
        "Future: hybrid search, cross-encoder re-ranking, RAGAS LLM-judge integration",
    ], footer_num="11/12")

    # 12. Thank You
    add_thank_you(prs)

    out = Path("docs/CiteMind_AI_Presentation.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f" Presentation saved → {out}")
    print(f"   {len(prs.slides)} slides")


if __name__ == "__main__":
    main()